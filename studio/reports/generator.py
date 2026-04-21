"""
LEWICE Studio — Certification Report Generator
Produces FAA-style PowerPoint (.pptx) reports with ice shapes,
conditions summary, impingement limits, and risk assessment.
"""
import os
import io
import tempfile
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    Image as RLImage,
    PageBreak,
)

# ── Brand Colors ──
NAVY       = RGBColor(0x0B, 0x1A, 0x2E)
DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)
GOLD       = RGBColor(0xC8, 0xA8, 0x4E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
TEXT_GRAY  = RGBColor(0x4A, 0x4A, 0x4A)
GREEN      = RGBColor(0x28, 0xA7, 0x45)
YELLOW     = RGBColor(0xE0, 0xA8, 0x00)
RED        = RGBColor(0xDC, 0x35, 0x45)

REPORTS_DIR = os.path.dirname(__file__)
STUDIO_DIR = os.path.abspath(os.path.join(REPORTS_DIR, ".."))
DEFAULT_LOGO_PATH = os.path.join(STUDIO_DIR, "assets", "tamarack-logo.png")


def _add_background(slide, color=NAVY):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_gold_line(slide, left, top, width):
    """Add a thin gold horizontal rule."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def _risk_color(risk_level):
    if risk_level == "LOW":
        return GREEN
    elif risk_level == "MEDIUM":
        return YELLOW
    else:
        return RED


def _assess_risk(max_thickness_mm, ice_type="unknown"):
    """Assess risk based on max ice thickness and ice type."""
    if max_thickness_mm < 5:
        return "LOW", "Thin accretion, minimal aerodynamic impact expected."
    elif max_thickness_mm < 15:
        return "MEDIUM", "Moderate accretion. Further analysis of aero impact recommended."
    else:
        return "HIGH", "Significant accretion with potential horn formations. Critical for certification."


def _classify_ice(temp_c, lwc, mvd):
    """Classify ice type based on conditions."""
    if temp_c < -15:
        return "Rime"
    elif temp_c > -5:
        return "Glaze"
    else:
        if lwc > 0.6:
            return "Mixed/Glaze"
        else:
            return "Mixed/Rime"


def generate_ice_shape_image(clean_coords, ice_coords, title="", width_px=900, height_px=500):
    """Generate an ice shape overlay plot as a PNG image bytes buffer.
    Uses matplotlib if available, otherwise returns None.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(1, 1, figsize=(width_px/100, height_px/100), dpi=100)
        fig.patch.set_facecolor('#0B1A2E')
        ax.set_facecolor('#0B1A2E')

        if clean_coords:
            cx = [p[0] for p in clean_coords]
            cy = [p[1] for p in clean_coords]
            ax.plot(cx, cy, color='#6699CC', linewidth=1.5, label='Clean Airfoil')

        if ice_coords:
            ix = [p[0] for p in ice_coords]
            iy = [p[1] for p in ice_coords]
            ax.plot(ix, iy, color='#C8A84E', linewidth=2.0, label='Iced Airfoil')

        ax.set_aspect('equal')
        ax.set_xlabel('x/c', color='#CCCCCC', fontsize=9)
        ax.set_ylabel('y/c', color='#CCCCCC', fontsize=9)
        ax.tick_params(colors='#888888', labelsize=8)
        ax.spines['bottom'].set_color('#444444')
        ax.spines['left'].set_color('#444444')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.legend(loc='upper right', fontsize=8, facecolor='#0B1A2E',
                  edgecolor='#444444', labelcolor='#CCCCCC')
        if title:
            ax.set_title(title, color='#C8A84E', fontsize=11, fontweight='bold', pad=10)
        ax.grid(True, alpha=0.15, color='#666666')

        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf
    except ImportError:
        return None


def generate_le_detail_image(clean_coords, ice_coords, title="Leading Edge Detail", width_px=900, height_px=500):
    """Generate a zoomed-in LE detail plot."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(1, 1, figsize=(width_px/100, height_px/100), dpi=100)
        fig.patch.set_facecolor('#0B1A2E')
        ax.set_facecolor('#0B1A2E')

        if clean_coords:
            cx = [p[0] for p in clean_coords]
            cy = [p[1] for p in clean_coords]
            ax.plot(cx, cy, color='#6699CC', linewidth=1.5, label='Clean Airfoil')

        if ice_coords:
            ix = [p[0] for p in ice_coords]
            iy = [p[1] for p in ice_coords]
            ax.plot(ix, iy, color='#C8A84E', linewidth=2.5, label='Iced Airfoil')
            ax.fill(ix, iy, color='#C8A84E', alpha=0.15)

        # Slightly wider window so LE shape is easier to read on slides.
        ax.set_xlim(-0.08, 0.30)
        ax.set_ylim(-0.12, 0.12)
        ax.set_aspect('equal')
        ax.set_xlabel('x/c', color='#CCCCCC', fontsize=9)
        ax.set_ylabel('y/c', color='#CCCCCC', fontsize=9)
        ax.tick_params(colors='#888888', labelsize=8)
        ax.spines['bottom'].set_color('#444444')
        ax.spines['left'].set_color('#444444')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.legend(loc='upper right', fontsize=8, facecolor='#0B1A2E',
                  edgecolor='#444444', labelcolor='#CCCCCC')
        ax.set_title(title, color='#C8A84E', fontsize=11, fontweight='bold', pad=10)
        ax.grid(True, alpha=0.15, color='#666666')

        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf
    except ImportError:
        return None


def generate_conditions_dashboard_image(conditions, max_thickness_mm, risk_level, width_px=700, height_px=500):
    """Generate a compact dashboard image for conditions and severity."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        speed = float(conditions.get("speed_ktas", 0) or 0)
        altitude = float(conditions.get("altitude_ft", 0) or 0)
        temp = float(conditions.get("temp_c", -10) or -10)
        lwc = float(conditions.get("lwc", 0.0) or 0.0)
        mvd = float(conditions.get("mvd", 0.0) or 0.0)

        # Normalize into 0..1 for an at-a-glance severity profile.
        metrics = {
            "Speed": min(max(speed / 300.0, 0.0), 1.0),
            "Altitude": min(max(altitude / 20000.0, 0.0), 1.0),
            "LWC": min(max(lwc / 1.2, 0.0), 1.0),
            "MVD": min(max(mvd / 40.0, 0.0), 1.0),
            "Ice Thick": min(max(max_thickness_mm / 25.0, 0.0), 1.0),
            "Warmth": min(max((temp + 30.0) / 30.0, 0.0), 1.0),
        }

        labels = list(metrics.keys())
        values = [metrics[k] for k in labels]

        fig = plt.figure(figsize=(width_px / 100, height_px / 100), dpi=100)
        fig.patch.set_facecolor("#0B1A2E")
        gs = fig.add_gridspec(2, 1, height_ratios=[3.0, 1.1], hspace=0.32)

        ax = fig.add_subplot(gs[0])
        ax.set_facecolor("#0B1A2E")
        bars = ax.barh(labels, values, color="#C8A84E", edgecolor="#E8E8E8", alpha=0.92)
        ax.set_xlim(0, 1.0)
        ax.invert_yaxis()
        ax.set_title("Icing Severity Profile", color="#C8A84E", fontsize=12, fontweight="bold", pad=10)
        ax.tick_params(colors="#CCCCCC", labelsize=9)
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.grid(axis="x", color="#445066", alpha=0.35)

        for bar, val in zip(bars, values):
            ax.text(min(val + 0.03, 0.98), bar.get_y() + bar.get_height() / 2, f"{val * 100:.0f}%",
                    va="center", ha="left", color="#FFFFFF", fontsize=8)

        ax2 = fig.add_subplot(gs[1])
        ax2.set_facecolor("#0B1A2E")
        ax2.axis("off")
        risk_color = {"LOW": "#28A745", "MEDIUM": "#E0A800", "HIGH": "#DC3545"}.get(risk_level, "#E0A800")
        summary = (
            f"Risk: {risk_level}   |   Temp: {temp:.1f} C   |   LWC: {lwc:.2f} g/m3   |   "
            f"MVD: {mvd:.0f} um   |   Max Thickness: {max_thickness_mm:.1f} mm"
        )
        ax2.text(0.02, 0.60, summary, color="#FFFFFF", fontsize=10, fontweight="bold", transform=ax2.transAxes)
        ax2.text(0.02, 0.20, "Higher bar values indicate more severe icing contributors for this scenario.",
                 color="#AEB8C4", fontsize=8.5, transform=ax2.transAxes)
        ax2.add_patch(plt.Rectangle((0.0, 0.02), 1.0, 0.04, transform=ax2.transAxes, color="#1A3A5C", ec="none"))
        ax2.add_patch(plt.Rectangle((0.0, 0.02), 0.25 if risk_level == "LOW" else 0.6 if risk_level == "MEDIUM" else 0.9,
                                    0.04, transform=ax2.transAxes, color=risk_color, ec="none"))

        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf
    except ImportError:
        return None


def build_report(
    project_name="Aircraft Icing Analysis",
    airfoil_name="NACA 0012",
    chord_m=0.9144,
    conditions=None,
    clean_coords=None,
    ice_coords=None,
    impingement_data=None,
    max_thickness_mm=0.0,
    analyst_name="",
    notes="",
):
    """Build a full certification-style PPTX report.

    Args:
        project_name: Name of the project/aircraft
        airfoil_name: Airfoil designation
        chord_m: Chord length in meters
        conditions: dict with keys: speed_ktas, altitude_ft, temp_c, lwc, mvd, exposure_min
        clean_coords: list of (x,y) for clean airfoil
        ice_coords: list of (x,y) for iced airfoil
        impingement_data: dict with upper/lower limits
        max_thickness_mm: maximum ice thickness in mm
        analyst_name: name of the analyst
        notes: additional notes

    Returns:
        BytesIO buffer containing the .pptx file
    """
    if conditions is None:
        conditions = {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    now = datetime.datetime.now().strftime("%B %d, %Y")
    temp_c = conditions.get("temp_c", -10)
    lwc = conditions.get("lwc", 0.5)
    mvd = conditions.get("mvd", 20)
    ice_type = _classify_ice(temp_c, lwc, mvd)
    risk_level, risk_desc = _assess_risk(max_thickness_mm, ice_type)
    chord_in = chord_m * 39.3700787402

    # ================================================================
    # SLIDE 1: Title
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, NAVY)

    # Gold accent bar top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    _add_textbox(slide, 1.0, 1.5, 11, 1.0, "LEWICE STUDIO", font_size=14,
                 color=GOLD, bold=True)
    _add_textbox(slide, 1.0, 2.0, 11, 1.5, "Ice Accretion Analysis Report",
                 font_size=36, color=WHITE, bold=True)
    _add_gold_line(slide, 1.0, 3.4, 4.0)
    _add_textbox(slide, 1.0, 3.7, 11, 0.5, project_name,
                 font_size=20, color=LIGHT_GRAY)
    _add_textbox(slide, 1.0, 4.3, 6.5, 0.4, f"Airfoil: {airfoil_name}  |  Chord: {chord_m:.3f} m ({chord_in:.2f} in)",
                 font_size=13, color=LIGHT_GRAY)
    _add_textbox(slide, 1.0, 4.8, 5, 0.4, f"Date: {now}",
                 font_size=12, color=TEXT_GRAY)
    if analyst_name:
        _add_textbox(slide, 1.0, 5.2, 5, 0.4, f"Analyst: {analyst_name}",
                     font_size=12, color=TEXT_GRAY)

    # Certification reference
    _add_textbox(slide, 1.0, 6.2, 11, 0.4,
                 "Reference: 14 CFR Part 25 Appendix C  |  NASA LEWICE 3.2  |  LEWICE Studio",
                 font_size=10, color=TEXT_GRAY)

    # Risk badge top-right
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.0), Inches(1.5), Inches(2.5), Inches(1.2)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = _risk_color(risk_level)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"RISK: {risk_level}"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # ================================================================
    # SLIDE 2: Conditions Summary
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, NAVY)

    _add_textbox(slide, 0.8, 0.4, 6, 0.6, "Icing Conditions Summary",
                 font_size=28, color=WHITE, bold=True)
    _add_gold_line(slide, 0.8, 1.0, 3.0)

    # Conditions table
    rows = [
        ("Parameter", "Value", "Unit"),
        ("Freestream Velocity", f"{conditions.get('speed_ktas', 'N/A')}", "KTAS"),
        ("Altitude", f"{conditions.get('altitude_ft', 'N/A')}", "ft MSL"),
        ("Static Temperature", f"{temp_c}", "deg C"),
        ("Liquid Water Content", f"{lwc}", "g/m^3"),
        ("Median Volume Diameter", f"{mvd}", "microns"),
        ("Exposure Time", f"{conditions.get('exposure_min', 'N/A')}", "min"),
        ("Angle of Attack", f"{conditions.get('aoa_deg', 'N/A')}", "deg"),
        ("Chord Length", f"{chord_m:.4f} / {chord_in:.2f}", "m / in"),
        ("Ice Type (Predicted)", ice_type, ""),
    ]

    table_shape = slide.shapes.add_table(len(rows), 3, Inches(0.8), Inches(1.3),
                                          Inches(6.5), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(1.5)

    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD
            else:
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE if i % 2 == 0 else NAVY

    # Right side: visual dashboard + key metrics
    dash_buf = generate_conditions_dashboard_image(conditions, max_thickness_mm, risk_level)
    _add_textbox(slide, 7.8, 1.1, 4.8, 0.4, "Condition Dashboard", font_size=14, color=GOLD, bold=True)
    if dash_buf:
        slide.shapes.add_picture(dash_buf, Inches(7.7), Inches(1.4), Inches(4.9), Inches(3.7))
    else:
        _add_textbox(slide, 7.9, 2.8, 4.4, 0.6, "Install matplotlib for dashboard visuals.",
                     font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, 8.0, 5.2, 4.5, 0.5, "Key Metrics", font_size=16, color=GOLD, bold=True)

    metrics = [
        ("Max Ice Thickness", f"{max_thickness_mm:.1f} mm"),
        ("Ice Type", ice_type),
        ("Risk Assessment", risk_level),
    ]
    y = 5.7
    for label, value in metrics:
        _add_textbox(slide, 8.0, y, 4.5, 0.3, label, font_size=10, color=TEXT_GRAY)
        _add_textbox(slide, 8.0, y + 0.22, 4.5, 0.38, value, font_size=16, color=GOLD, bold=True)
        y += 0.58

    # ================================================================
    # SLIDE 3: Ice Shape — Full Airfoil
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, NAVY)

    _add_textbox(slide, 0.8, 0.4, 6, 0.6, "Ice Shape — Full Airfoil View",
                 font_size=28, color=WHITE, bold=True)
    _add_gold_line(slide, 0.8, 1.0, 3.0)

    img_buf = generate_ice_shape_image(
        clean_coords, ice_coords,
        title=f"{airfoil_name} — {ice_type} Ice at {temp_c} C, {lwc} g/m3, {mvd} um MVD"
    )
    if img_buf:
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5))
    else:
        _add_textbox(slide, 2, 3, 8, 1, "[Ice shape plot — install matplotlib for images]",
                     font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 4: Ice Shape — LE Detail
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, NAVY)

    _add_textbox(slide, 0.8, 0.4, 6, 0.6, "Ice Shape — Leading Edge Detail",
                 font_size=28, color=WHITE, bold=True)
    _add_gold_line(slide, 0.8, 1.0, 3.0)

    le_buf = generate_le_detail_image(
        clean_coords, ice_coords,
        title=f"LE Detail — Max Thickness: {max_thickness_mm:.1f} mm"
    )
    if le_buf:
        slide.shapes.add_picture(le_buf, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5))
    else:
        _add_textbox(slide, 2, 3, 8, 1, "[LE detail plot — install matplotlib for images]",
                     font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 5: Impingement Limits & Ice Metrics
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, NAVY)

    _add_textbox(slide, 0.8, 0.4, 6, 0.6, "Impingement Limits & Ice Metrics",
                 font_size=28, color=WHITE, bold=True)
    _add_gold_line(slide, 0.8, 1.0, 3.0)

    if impingement_data:
        imp_rows = [("Metric", "x/c", "y/c", "s/c")]
        if "upper" in impingement_data:
            u = impingement_data["upper"]
            imp_rows.append(("Upper Limit", f"{u.get('x', 'N/A'):.6f}",
                           f"{u.get('y', 'N/A'):.6f}", f"{u.get('s', 'N/A'):.6f}"))
        if "lower" in impingement_data:
            lo = impingement_data["lower"]
            imp_rows.append(("Lower Limit", f"{lo.get('x', 'N/A'):.6f}",
                           f"{lo.get('y', 'N/A'):.6f}", f"{lo.get('s', 'N/A'):.6f}"))
        if "stagnation" in impingement_data:
            st = impingement_data["stagnation"]
            imp_rows.append(("Stagnation", f"{st.get('x', 'N/A'):.6f}",
                           f"{st.get('y', 'N/A'):.6f}", "—"))

        table_shape = slide.shapes.add_table(len(imp_rows), 4, Inches(0.8), Inches(1.4),
                                              Inches(7.0), Inches(1.5))
        table = table_shape.table
        for i, row_data in enumerate(imp_rows):
            for j, val in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.name = "Calibri"
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = NAVY
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GOLD
                else:
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = DARK_BLUE if i % 2 == 0 else NAVY
    else:
        _add_textbox(slide, 0.8, 1.5, 8, 0.5, "Impingement data not available for this run.",
                     font_size=13, color=TEXT_GRAY)

    # Ice metrics summary
    _add_textbox(slide, 0.8, 3.5, 6, 0.5, "Accretion Summary", font_size=18, color=GOLD, bold=True)

    summary_items = [
        f"Maximum ice thickness:  {max_thickness_mm:.1f} mm  ({max_thickness_mm/25.4:.3f} in)",
        f"Ice type classification:  {ice_type}",
        f"Risk assessment:  {risk_level} — {risk_desc}",
        f"Icing exposure:  {conditions.get('exposure_min', 'N/A')} min at {lwc} g/m3 LWC",
        f"Clean airfoil points:  {len(clean_coords) if clean_coords else 'N/A'}",
        f"Iced airfoil points:  {len(ice_coords) if ice_coords else 'N/A'}",
    ]
    for idx, item in enumerate(summary_items):
        _add_textbox(slide, 1.0, 4.1 + idx * 0.4, 10, 0.35, item, font_size=12, color=WHITE)

    # ================================================================
    # SLIDE 6: Risk Assessment & Recommendations
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, NAVY)

    _add_textbox(slide, 0.8, 0.4, 6, 0.6, "Risk Assessment & Recommendations",
                 font_size=28, color=WHITE, bold=True)
    _add_gold_line(slide, 0.8, 1.0, 3.0)

    # Traffic light
    lights = [("LOW", GREEN), ("MEDIUM", YELLOW), ("HIGH", RED)]
    for idx, (label, color) in enumerate(lights):
        x = 1.5 + idx * 3.5
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.5), Inches(1.2), Inches(1.2))
        if label == risk_level:
            circ.fill.solid()
            circ.fill.fore_color.rgb = color
        else:
            circ.fill.solid()
            circ.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
        circ.line.color.rgb = color
        circ.line.width = Pt(2)

        _add_textbox(slide, x - 0.15, 2.8, 1.5, 0.4, label,
                     font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Risk description
    _add_textbox(slide, 0.8, 3.5, 10, 0.5, f"Assessment: {risk_level}",
                 font_size=22, color=_risk_color(risk_level), bold=True)
    _add_textbox(slide, 0.8, 4.1, 10, 0.5, risk_desc, font_size=14, color=WHITE)

    # Recommendations based on risk
    _add_textbox(slide, 0.8, 4.8, 6, 0.5, "Recommendations", font_size=18, color=GOLD, bold=True)

    if risk_level == "LOW":
        recs = [
            "Ice accretion is within acceptable limits for continued flight operations.",
            "Standard ice protection system activation procedures are sufficient.",
            "No additional aerodynamic penalty analysis required for this condition.",
            "Document results for certification file."
        ]
    elif risk_level == "MEDIUM":
        recs = [
            "Conduct aerodynamic performance degradation analysis (Cl_max, Cd penalty).",
            "Verify ice protection system coverage extends to impingement limits.",
            "Consider additional LEWICE runs at extended exposure times.",
            "Review against critical ice shape criteria per AC 20-73A.",
            "Document results and submit for DER/ODA review."
        ]
    else:
        recs = [
            "CRITICAL: Significant ice accretion with potential horn formation detected.",
            "Full aerodynamic analysis required (CFD or wind tunnel with ice shape).",
            "Verify ice protection system can handle this accretion rate.",
            "Run parametric study across full Appendix C envelope.",
            "Consider ice shape fidelity testing per AC 20-73A Section 7.",
            "Submit to DER/ODA for review prior to certification credit."
        ]

    for idx, rec in enumerate(recs):
        prefix = "+" if risk_level == "LOW" else "!" if risk_level == "HIGH" else ">"
        _add_textbox(slide, 1.0, 5.3 + idx * 0.35, 11, 0.3,
                     f"  {prefix}  {rec}", font_size=11, color=WHITE)

    # Notes
    if notes:
        _add_textbox(slide, 0.8, 6.8, 10, 0.4, f"Notes: {notes}", font_size=10, color=TEXT_GRAY)

    # Footer on all slides
    for slide in prs.slides:
        _add_textbox(slide, 0.5, 7.0, 12, 0.35,
                     f"LEWICE Studio  |  {project_name}  |  {now}  |  Confidential",
                     font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Save to buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _chord_units(chord_m):
    chord_in = chord_m * 39.3700787402
    return chord_m, chord_in


def build_pdf_report(
    project_name="Aircraft Icing Analysis",
    airfoil_name="NACA 0012",
    chord_m=0.9144,
    conditions=None,
    clean_coords=None,
    ice_coords=None,
    impingement_data=None,
    max_thickness_mm=0.0,
    analyst_name="",
    notes="",
    logo_path=DEFAULT_LOGO_PATH,
):
    """Build a certification-style PDF report for FAA icing analysis workflows."""
    if conditions is None:
        conditions = {}

    now = datetime.datetime.now().strftime("%B %d, %Y")
    temp_c = float(conditions.get("temp_c", -10) or -10)
    lwc = float(conditions.get("lwc", 0.5) or 0.5)
    mvd = float(conditions.get("mvd", 20) or 20)
    ice_type = _classify_ice(temp_c, lwc, mvd)
    risk_level, risk_desc = _assess_risk(max_thickness_mm, ice_type)
    chord_m_val, chord_in = _chord_units(chord_m)

    risk_rgb = {
        "LOW": colors.HexColor("#28A745"),
        "MEDIUM": colors.HexColor("#E0A800"),
        "HIGH": colors.HexColor("#DC3545"),
    }.get(risk_level, colors.HexColor("#E0A800"))
    risk_hex = {
        "LOW": "#28A745",
        "MEDIUM": "#E0A800",
        "HIGH": "#DC3545",
    }.get(risk_level, "#E0A800")

    regulations = [
        "14 CFR Part 25, Appendix C - Icing Envelope (continuous maximum and intermittent maximum).",
        "14 CFR 25.1419 - Ice protection system capability for dispatch and operation in icing.",
        "FAA AC 20-73A - Aircraft Ice Protection guidance for analysis and substantiation practices.",
        "NASA LEWICE User/Validation references for computational ice accretion prediction methods.",
    ]

    toc_sections = [
        "1. Executive Summary",
        "2. Applicable Regulations and Guidance",
        "3. LEWICE Theory and Method Basis",
        "4. Icing Conditions and Inputs",
        "5. Ice Shape Results (Full Airfoil and LE Detail)",
        "6. Impingement Limits and Ice Metrics",
        "7. Risk Assessment and Recommendations",
        "8. Analyst Notes",
    ]

    theory_paras = [
        "LEWICE computes droplet trajectories, collection efficiency, and thermodynamic freezing behavior along the airfoil surface to estimate accreted ice geometry.",
        "The solver combines airflow, droplet impingement physics, and an energy balance to resolve whether incoming liquid water freezes, runs back, or sheds.",
        "Predicted geometry (for example horn or rime shapes) is then used as a certification input for aerodynamic impact and ice protection capability assessment.",
    ]

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=letter,
        leftMargin=0.6 * inch,
        rightMargin=0.6 * inch,
        topMargin=1.0 * inch,
        bottomMargin=0.65 * inch,
        title=f"{project_name} - LEWICE Certification Report",
        author=analyst_name or "LEWICE Studio",
    )

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=17, leading=21, textColor=colors.HexColor("#0B1A2E"), spaceAfter=8)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=13, leading=17, textColor=colors.HexColor("#1A3A5C"), spaceAfter=6)
    body = ParagraphStyle("Body", parent=styles["BodyText"], fontSize=10, leading=13, spaceAfter=6)
    small = ParagraphStyle("Small", parent=styles["BodyText"], fontSize=8.5, leading=11, textColor=colors.HexColor("#4A4A4A"))

    story = []

    story.append(Paragraph("FAA-Style Icing Certification Report", h1))
    story.append(Paragraph(f"<b>Project:</b> {project_name}", body))
    story.append(Paragraph(f"<b>Airfoil:</b> {airfoil_name}", body))
    story.append(Paragraph(f"<b>Chord:</b> {chord_m_val:.4f} m ({chord_in:.2f} in)", body))
    story.append(Paragraph(f"<b>Date:</b> {now}", body))
    if analyst_name:
        story.append(Paragraph(f"<b>Analyst:</b> {analyst_name}", body))
    story.append(Spacer(1, 0.1 * inch))
    story.append(Paragraph(f"<b>Executive Risk:</b> <font color='{risk_hex}'>{risk_level}</font>", body))
    story.append(Paragraph(risk_desc, body))
    story.append(PageBreak())

    story.append(Paragraph("Table of Contents", h1))
    for item in toc_sections:
        story.append(Paragraph(item, body))
    story.append(PageBreak())

    story.append(Paragraph("1. Executive Summary", h1))
    story.append(Paragraph(
        "This report summarizes NASA LEWICE-generated ice accretion geometry and certification-relevant metrics for the stated flight and icing conditions.",
        body,
    ))
    story.append(Paragraph(
        f"Predicted ice type is <b>{ice_type}</b> with maximum estimated leading-edge accretion of <b>{max_thickness_mm:.1f} mm ({max_thickness_mm/25.4:.3f} in)</b>.",
        body,
    ))

    story.append(Paragraph("2. Applicable Regulations and Guidance", h1))
    for line in regulations:
        story.append(Paragraph(f"- {line}", body))

    story.append(Paragraph("3. LEWICE Theory and Method Basis", h1))
    for para in theory_paras:
        story.append(Paragraph(para, body))

    story.append(PageBreak())
    story.append(Paragraph("4. Icing Conditions and Inputs", h1))
    cond_rows = [
        ["Parameter", "Value", "Unit"],
        ["Freestream Velocity", f"{conditions.get('speed_ktas', 'N/A')}", "KTAS"],
        ["Altitude", f"{conditions.get('altitude_ft', 'N/A')}", "ft MSL"],
        ["Static Temperature", f"{temp_c:.2f}", "deg C"],
        ["Liquid Water Content", f"{lwc:.3f}", "g/m^3"],
        ["Median Volume Diameter", f"{mvd:.1f}", "microns"],
        ["Exposure Time", f"{conditions.get('exposure_min', 'N/A')}", "min"],
        ["Angle of Attack", f"{conditions.get('aoa_deg', 'N/A')}", "deg"],
        ["Chord Length", f"{chord_m_val:.4f} / {chord_in:.2f}", "m / in"],
    ]
    cond_table = Table(cond_rows, colWidths=[2.8 * inch, 2.3 * inch, 1.4 * inch])
    cond_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#C8A84E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0B1A2E")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("BACKGROUND", (0, 1), (-1, -1), colors.whitesmoke),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#B0B0B0")),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story.append(cond_table)
    story.append(Spacer(1, 0.12 * inch))
    story.append(Paragraph(f"Predicted ice type: <b>{ice_type}</b>", body))

    story.append(Paragraph("5. Ice Shape Results (Full Airfoil and LE Detail)", h1))
    full_buf = generate_ice_shape_image(
        clean_coords,
        ice_coords,
        title=f"{airfoil_name} - Full Airfoil",
        width_px=1200,
        height_px=640,
    )
    if full_buf:
        story.append(RLImage(full_buf, width=6.9 * inch))
    else:
        story.append(Paragraph("Matplotlib not available: full-airfoil plot could not be rendered.", body))

    le_buf = generate_le_detail_image(
        clean_coords,
        ice_coords,
        title=f"Leading Edge Detail - {max_thickness_mm:.1f} mm",
        width_px=1200,
        height_px=640,
    )
    if le_buf:
        story.append(Spacer(1, 0.08 * inch))
        story.append(RLImage(le_buf, width=6.9 * inch, height=3.5 * inch))
    else:
        story.append(Paragraph("Matplotlib not available: LE detail plot could not be rendered.", body))

    story.append(PageBreak())
    story.append(Paragraph("6. Impingement Limits and Ice Metrics", h1))
    if impingement_data:
        imp_rows = [["Metric", "x/c", "y/c", "s/c"]]
        if "upper" in impingement_data:
            u = impingement_data["upper"]
            imp_rows.append(["Upper Limit", f"{u.get('x', 0):.6f}", f"{u.get('y', 0):.6f}", f"{u.get('s', 0):.6f}"])
        if "lower" in impingement_data:
            lo = impingement_data["lower"]
            imp_rows.append(["Lower Limit", f"{lo.get('x', 0):.6f}", f"{lo.get('y', 0):.6f}", f"{lo.get('s', 0):.6f}"])
        if "stagnation" in impingement_data:
            stg = impingement_data["stagnation"]
            imp_rows.append(["Stagnation", f"{stg.get('x', 0):.6f}", f"{stg.get('y', 0):.6f}", "-"])

        imp_table = Table(imp_rows, colWidths=[2.0 * inch, 1.5 * inch, 1.5 * inch, 1.5 * inch])
        imp_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#C8A84E")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0B1A2E")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#B0B0B0")),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        story.append(imp_table)
        story.append(Spacer(1, 0.12 * inch))
    else:
        story.append(Paragraph("Impingement limit data not available for this run.", body))

    story.append(Paragraph(f"Maximum ice thickness: <b>{max_thickness_mm:.1f} mm ({max_thickness_mm/25.4:.3f} in)</b>", body))
    story.append(Paragraph(f"Risk assessment: <font color='{risk_hex}'><b>{risk_level}</b></font> - {risk_desc}", body))
    story.append(Paragraph(f"Clean airfoil points: {len(clean_coords) if clean_coords else 'N/A'}", body))
    story.append(Paragraph(f"Iced airfoil points: {len(ice_coords) if ice_coords else 'N/A'}", body))

    story.append(Paragraph("7. Risk Assessment and Recommendations", h1))
    if risk_level == "LOW":
        recs = [
            "Ice accretion is within acceptable limits for baseline operation.",
            "Standard ice protection activation procedures are likely sufficient.",
            "Archive this run in certification records as supporting evidence.",
        ]
    elif risk_level == "MEDIUM":
        recs = [
            "Perform aerodynamic penalty assessment (stall margin and drag rise).",
            "Verify protected surface coverage against impingement limits.",
            "Run additional exposure-time sensitivity points for margin.",
        ]
    else:
        recs = [
            "Critical accretion detected; complete high-fidelity aerodynamic substantiation.",
            "Expand parametric runs across Appendix C critical combinations.",
            "Present results for DER/ODA review prior to certification credit request.",
        ]
    for rec in recs:
        story.append(Paragraph(f"- {rec}", body))

    story.append(Paragraph("8. Analyst Notes", h1))
    story.append(Paragraph(notes if notes else "No additional analyst notes provided.", body))
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph(
        "Generated by LEWICE Studio using NASA LEWICE 3.2 computational output for certification-oriented reporting.",
        small,
    ))

    def draw_header_footer(canvas, document):
        canvas.saveState()
        page_w, page_h = letter

        canvas.setFillColor(colors.HexColor("#0B1A2E"))
        canvas.rect(0, page_h - 0.62 * inch, page_w, 0.62 * inch, stroke=0, fill=1)
        canvas.setFillColor(colors.HexColor("#C8A84E"))
        canvas.rect(0, page_h - 0.64 * inch, page_w, 0.02 * inch, stroke=0, fill=1)

        if logo_path and os.path.isfile(logo_path):
            canvas.drawImage(ImageReader(logo_path), 0.62 * inch, page_h - 0.54 * inch, width=1.05 * inch, height=0.34 * inch, preserveAspectRatio=True, mask='auto')

        canvas.setFillColor(colors.white)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawRightString(page_w - 0.62 * inch, page_h - 0.36 * inch, "Tamarack Aerospace - Icing Certification Report")

        canvas.setFillColor(colors.HexColor("#4A4A4A"))
        canvas.setFont("Helvetica", 8)
        canvas.drawString(0.62 * inch, 0.42 * inch, f"LEWICE Studio | {project_name} | {now}")
        canvas.drawRightString(page_w - 0.62 * inch, 0.42 * inch, f"Page {document.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_header_footer, onLaterPages=draw_header_footer)
    buf.seek(0)
    return buf


def save_report(filepath, **kwargs):
    """Save report to a .pptx file on disk."""
    buf = build_report(**kwargs)
    with open(filepath, "wb") as f:
        f.write(buf.read())
    return filepath


if __name__ == "__main__":
    import sys
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
    from lewice_engine.output_parser import parse_clean_airfoil, parse_ice_shape

    clean = parse_clean_airfoil("../../Inputs/case1.xyd")
    ice = parse_ice_shape("../../studio/test_run/case1/final1.dat")

    min_x_ice = min(c[0] for c in ice)
    max_thickness = abs(min_x_ice) * 0.9144 * 1000

    out = save_report(
        "test_report.pptx",
        project_name="Case 1 — NASA LEWICE Validation",
        airfoil_name="NACA 0012 (modified)",
        chord_m=0.9144,
        conditions={
            "speed_ktas": 175,
            "altitude_ft": 10000,
            "temp_c": -4.85,
            "lwc": 0.54,
            "mvd": 20,
            "exposure_min": 6.0,
            "aoa_deg": 4.5,
        },
        clean_coords=clean,
        ice_coords=ice,
        impingement_data={
            "upper": {"x": 0.003625, "y": 0.005619, "s": 0.020414},
            "lower": {"x": 0.145164, "y": -0.030770, "s": -0.143421},
        },
        max_thickness_mm=max_thickness,
        analyst_name="LEWICE Studio",
        notes="NASA LEWICE 3.2 Case 1 validation run.",
    )
    print(f"Report saved: {out}")
